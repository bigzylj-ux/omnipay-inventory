#!/usr/bin/env python3
"""
Generate comprehensive OmniPay POS Inventory System PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(16, 120, 80)  # Emerald green
SECONDARY_COLOR = RGBColor(71, 85, 119)  # Dark blue
ACCENT_COLOR = RGBColor(244, 245, 247)  # Light gray
TEXT_COLOR = RGBColor(0, 0, 0)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title background
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(8.5), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title background
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    return slide

# ==================== SLIDE 1: TITLE ====================
add_title_slide(prs, 'OmniPay', 'POS Inventory System\nComprehensive Management Guide')

# ==================== SLIDE 2: EXECUTIVE SUMMARY ====================
add_content_slide(prs, 'Executive Summary', [
    '• Shared POS terminal inventory workflow for admins and approved users',
    '• Supabase-first storage with IndexedDB fallback for resilience',
    '• Real-time dashboard with KPI tracking and charts',
    '• Daily automated reconciliation of inventory vs portal data',
    '• Vendor management, repair cost tracking, and fault uploads',
    '• Large-upload support for thousands of records through chunked processing',
])

# ==================== SLIDE 3: KEY FEATURES ====================
add_content_slide(prs, 'Key Features', [
    '✓ Master Inventory Management - Track all terminals with detailed metadata',
    '✓ Daily Portal Reconciliation - Auto-match and update terminal assignments',
    '✓ Vendor Management - Manage repair vendors and fault-based costs',
    '✓ Advanced Search & Filtering - By serial, terminal ID, merchant, location, status',
    '✓ Real-time Dashboard - KPI cards showing inventory status overview',
    '✓ Export Capabilities - Excel export for reporting and backup',
    '✓ Audit Trail - Complete reconciliation history and logs',
])

# ==================== SLIDE 4: SYSTEM ARCHITECTURE ====================
add_content_slide(prs, 'System Architecture', [
    'User Interface Layer',
    '  → React 18.2 + TypeScript for the inventory workflow UI',
    '',
    'Auth & Access Layer',
    '  → Protected routes, admin approvals, and role-based access',
    '',
    'Data Layer',
    '  → Supabase for shared storage, with IndexedDB fallback when needed',
    '',
    'Operations Layer',
    '  → Excel import, reconciliation, tracking, and audit logging',
])

# ==================== SLIDE 5: TECHNOLOGY STACK ====================
add_two_column_slide(prs, 'Technology Stack', [
    'Frontend:',
    '• React 18.2',
    '• TypeScript 5.3',
    '• React Router v6',
    '• Tailwind CSS 3.3',
    '• Lucide React (Icons)',
], [
    'Build & Tools:',
    '• Vite 5.0 (Fast bundler)',
    '• PostCSS 8.4',
    '• Node.js / npm',
    '',
    'Data:',
    '• IndexedDB (Browser)',
    '• Excel/CSV import',
])

# ==================== SLIDE 6: CORE COMPONENTS ====================
add_content_slide(prs, 'Core Components', [
    '📊 Dashboard - Real-time KPI overview of inventory status',
    '📦 Inventory - Searchable master inventory with advanced filtering',
    '📥 Daily Import - Master data and portal DB upload with deduplication',
    '✔️ Reconciliation - Automated matching and auto-update engine',
    '🏢 Vendors - Vendor contact management and repair tracking',
    '🗺️ Tracking - Historical audit trail of terminal movements',
])

# ==================== SLIDE 7: DAILY WORKFLOW ====================
add_content_slide(prs, 'Daily Workflow', [
    '1. Download portal DB file from your portal system',
    '',
    '2. Go to Daily Import → Portal DB tab',
    '',
    '3. Upload the DB file (auto-deduplication & validation)',
    '',
    '4. Navigate to Reconciliation page',
    '',
    '5. Click "Run Daily Reconciliation"',
    '',
    '6. Review results (updated, new, no change, exceptions)',
])

# ==================== SLIDE 8: MASTER INVENTORY ====================
add_content_slide(prs, 'Master Inventory Setup', [
    'One-time Initial Setup:',
    '',
    '1. Download Master Inventory Template from Daily Import page',
    '',
    '2. Prepare your inventory data with columns:',
    '   SN, Device Serial No., TerminalID, MerchantName, Status, etc.',
    '',
    '3. Upload via Daily Import → Master Inventory tab',
    '',
    '4. System imports all terminals into local database',
    '',
    'After Setup: Only re-upload when adding new terminal batches',
])

# ==================== SLIDE 9: RECONCILIATION ENGINE ====================
add_content_slide(prs, 'Reconciliation Engine', [
    'Automatic Matching:',
    '• Matches portal serial numbers against master inventory',
    '• Auto-assigns terminal IDs where available',
    '• Updates merchant names from portal data',
    '• Updates phone numbers from portal',
    '',
    'Results Display:',
    '• Records Updated - Successfully matched and updated',
    '• New Deployments - New terminals found in portal',
    '• No Change - No updates needed',
    '• Exceptions Flagged - Mismatches requiring manual review',
])

# ==================== SLIDE 10: VENDOR MANAGEMENT ====================
add_content_slide(prs, 'Vendor Management', [
    'Add & Manage Vendors:',
    '• Store vendor contact details (email, phone, address)',
    '• Upload repair records with fault-based costs',
    '• Track repair history per terminal',
    '',
    'Update Terminal Status:',
    '• Mark as "Repaired" when repair completes',
    '• Mark as "Cannibalised" for salvage operations',
    '• Automatically update inventory status',
])

# ==================== SLIDE 11: SEARCH & FILTERING ====================
add_content_slide(prs, 'Advanced Search & Filtering', [
    'Search Across:',
    '• Device Serial Number',
    '• Terminal ID',
    '• Merchant Name',
    '• Phone Number',
    '',
    'Filter By:',
    '• Status (All statuses available)',
    '• Location',
    '• Category',
    '',
    'Action: Click "Export" to download filtered results as Excel',
])

# ==================== SLIDE 12: TERMINAL STATUS TYPES ====================
add_content_slide(prs, 'Terminal Status Types', [
    '• Deployed - Actively in use at merchant location',
    '• In Stock - Available in warehouse',
    '• Yet To Deploy - Ready for deployment',
    '• Under Repair - Sent to vendor for repair',
    '• Repaired - Repair completed successfully',
    '• Faulty - Terminal has identified defects',
    '• Retrieved - Recovered from field',
    '• Cannibalised - Parts salvaged',
    '• Lost - Terminal missing',
    '• Test Terminal - Used for testing purposes',
])

# ==================== SLIDE 13: TROUBLESHOOTING ====================
add_two_column_slide(prs, 'Troubleshooting Guide', [
    'Import Issues:',
    '• Missing sheet names',
    '• Wrong column headers',
    '• Solution: Use template',
    '',
    'No Portal Records:',
    '• DB file not imported',
    '• Solution: Import DB first',
    '',
    'Data Disappears:',
    '• Browser cache cleared',
    '• Solution: Export backups',
], [
    'Search Not Working:',
    '• No data imported yet',
    '• Filter mismatch',
    '• Solution: Import first',
    '',
    'Exceptions Flagged:',
    '• Serial number mismatch',
    '• Missing records',
    '• Solution: Manual review',
    '',
    'Performance Issues:',
    '• Large dataset (>10K)',
    '• Solution: Use filtering',
])

# ==================== SLIDE 14: DATA STORAGE & BACKUP ====================
add_content_slide(prs, 'Data Storage & Backup Strategy', [
    'Data Location:',
    '• Stored through Supabase when available, with IndexedDB fallback',
    '• Supports shared visibility for authorized users across browsers',
    '• Keeps audit history and operational records available for review',
    '',
    'Backup Best Practices:',
    '• Export inventory and logs regularly to Excel',
    '• Keep source import files for traceability',
    '• Archive historical records for audit and reporting',
    '• Validate shared access settings for team workflows',
])

# ==================== SLIDE 15: DAILY CHECKLIST ====================
add_content_slide(prs, 'Daily Operations Checklist', [
    '☑ Download portal DB sheet from portal system',
    '☑ Open OmniPay system',
    '☑ Navigate to Daily Import page',
    '☑ Upload portal DB file',
    '☑ Go to Reconciliation page',
    '☑ Run Daily Reconciliation',
    '☑ Review results (updated, new, exceptions)',
    '☑ Export reconciliation logs if needed',
])

# ==================== SLIDE 16: SYSTEM REQUIREMENTS ====================
add_content_slide(prs, 'System Requirements', [
    'Browser:',
    '• Modern web browser (Chrome, Firefox, Edge, Safari)',
    '• Recent versions recommended',
    '',
    'Settings:',
    '• JavaScript enabled',
    '• Cookies/Local Storage enabled',
    '• Pop-ups allowed for downloads',
    '',
    'Hardware:',
    '• Minimum 50MB disk space for full database',
    '• Stable internet for initial access',
    '• Works offline after loading',
])

# ==================== SLIDE 17: SECURITY & DATA RETENTION ====================
add_content_slide(prs, 'Security & Data Retention', [
    'Data Retention:',
    '• Master Inventory: Persists in shared storage and fallback storage',
    '• Portal Records: Latest import stored for current workflows',
    '• Reconciliation Logs: Retained for review and audit tracking',
    '• Vendor Data: Maintained for operational history',
    '',
    'Security Notes:',
    '• Role-based access controls protect key workflows',
    '• Shared storage is managed through Supabase policies',
    '• Regular exports and validation are recommended for compliance',
])

# ==================== SLIDE 18: IMPLEMENTATION HIGHLIGHTS ====================
add_content_slide(prs, 'Implementation Highlights', [
    '✓ Intelligent Portal Deduplication',
    '  Removes duplicate records, keeps most complete entry',
    '',
    '✓ Smart Status Normalization',
    '  Auto-maps various status values to standard statuses',
    '',
    '✓ Exception Flagging System',
    '  Alerts on mismatches (mismatched Terminal IDs, missing records)',
    '',
    '✓ Fault-Based Cost Tracking',
    '  Tracks repair costs by fault type per vendor',
])

# ==================== SLIDE 19: FUTURE ROADMAP ====================
add_content_slide(prs, 'Future Roadmap & Enhancements', [
    'Phase 2 - Advanced Operations:',
    '• Expanded analytics and forecasting reports',
    '• More automation for exception review workflows',
    '• Stronger reporting for regional and vendor performance',
    '',
    'Phase 3 - Integrations:',
    '• Mobile support for field tracking',
    '• Automated notifications and reminders',
    '• Integration with downstream finance and audit tools',
])

# ==================== SLIDE 20: FAQ ====================
add_content_slide(prs, 'Frequently Asked Questions', [
    'Q: Do I need internet always?',
    'A: The app is best used with internet access for shared Supabase data, but fallback storage helps when needed.',
    '',
    'Q: Where is my data stored?',
    'A: The app prefers Supabase and uses IndexedDB fallback when the shared storage layer is unavailable.',
    '',
    'Q: Can multiple people access it?',
    'A: Yes, authorized users can access shared records when the correct permissions and Supabase setup are in place.',
    '',
    'Q: What if I clear browser cache?',
    'A: Browser-only copies may be lost, so regular exports and shared storage checks are recommended.',
])

# ==================== SLIDE 21: GETTING STARTED ====================
add_content_slide(prs, 'Getting Started - New Employee', [
    '1. Read Executive Summary (this presentation)',
    '',
    '2. Review Daily Workflow & Operations Checklist',
    '',
    '3. Practice Master Inventory upload with template',
    '',
    '4. Perform a test reconciliation',
    '',
    '5. Familiarize with search and filtering features',
    '',
    '6. Keep detailed documentation nearby for reference',
])

# ==================== SLIDE 22: CONTACT & SUPPORT ====================
add_title_slide(prs, 'Questions?', 'Refer to the complete documentation guide\nfor detailed procedures and troubleshooting')

# Save presentation
output_path = r'c:\Users\Sunday Aderibigbe\omnipay-inventory\OmniPay_System_Presentation.pptx'
prs.save(output_path)
print(f"✓ PowerPoint presentation generated successfully!")
print(f"✓ File saved to: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
